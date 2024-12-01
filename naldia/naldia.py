import random
import fitz
import docx
import json
import jsonlines
from tqdm import tqdm
import time
import validators
from bs4 import BeautifulSoup
import requests
import openai
import numpy as np
from numpy.linalg import norm
import os
import hashlib
import tiktoken
from pptx import Presentation
 

# Utilisez os.path.expandvars pour accéder à la variable d'environnement utilisateur
# nom_variable = "OPENAI_API_KEY"
# valeur_variable = os.path.expandvars(f"%{nom_variable}%")

if openai.api_key == "sk-" or openai.api_key == "" :
    print("\033[91m{}\033[0m".format("INVALID OPENAI API KEY"))
else:
    print("Clef d'API chargée.")


# openai.organization = ""

# definition des apî => développer diag classes Naldia, implémenter une factory
tokenizer = tiktoken.encoding_for_model("gpt-3.5-turbo")


# Namespace Naldia
class Naldia:
    
    offline = False
    async_mode = False

    api_model_ask = "text-davinci-003"
    api_model_memorize = "gpt-3.5-turbo"
    api_model_current = api_model_memorize
    
    api_call_wait = 0
    

    tokenizer = tiktoken.encoding_for_model(api_model_current)
  
    def sleep(seconds="api_call_wait"):
        if not Naldia.offline:
            if seconds == "api_call_wait" :
                time.sleep(Naldia.api_call_wait)
            else:
                time.sleep(seconds)
        return True
        
    def set_model(model):
        Naldia.api_model_current = model
        Naldia.tokenizer = tiktoken.encoding_for_model("gpt-3.5-turbo")
        Naldia.set_call_wait()
    
    def set_call_wait():
        if Naldia.api_model_current == Naldia.api_model_ask:
            Naldia.api_call_wait = 0
        elif Naldia.api_model_current == Naldia.api_model_memorize :
            Naldia.api_call_wait = 0

    # basic usage of openai API
    def ask(messages):
        Naldia.set_model(Naldia.api_model_current)
        if Naldia.api_model_current == Naldia.api_model_ask:
            response = Naldia.OpenAI_Completion(messages)
        elif Naldia.api_model_current == Naldia.api_model_memorize :
            response = Naldia.OpenAI_ChatCompletion(messages)
        return response

    def OpenAI_Completion(messages):
        # par analogie avec ChatCompletion, on reçois un message en entrée = [{"role": "user", "content": content}]
        # xuxa
        prompt = messages[0]["content"]
        max_tokens = 4096 - 3900
        try:
            # Utilisez la réponse de l'API comme souhaité
            response = openai.Completion.create(
                engine = Naldia.api_model_ask,
                prompt=prompt,
                temperature=1,
                max_tokens=max_tokens
            )
        except openai.OpenAIError as e:
            # Gérez les erreurs spécifiques de l'API OpenAI
            response = f"Une erreur OpenAI s'est produite : {e}"
            print(response)
        except Exception as e:
            # Gérez les autres exceptions inattendues
            print(f"Une erreur inattendue s'est produite : {e}")

        return response.choices[0].text.strip()

    def OpenAI_ChatCompletion(messages):
        try:
            if Naldia.async_mode:
                return Naldia.aChatCompletion(messages)
            else:
                return Naldia.ChatCompletion(messages)
            
        except openai.OpenAIError as e:
            response = f"Une erreur OpenAI s'est produite : {e}"
            print(response)
        return response

    def ChatCompletion(messages):
        try:
            completion = openai.ChatCompletion.create(
                model=Naldia.api_model_memorize,
                messages=messages,
                temperature=1,
                top_p=0.95,
                frequency_penalty=0.0,
                presence_penalty=0.0,
            )
            response = completion.choices[0].message.content
        except openai.OpenAIError as e:
            response = f"Une erreur OpenAI s'est produite : {e}"
            print(response)
        return response
    
    async def aChatCompletion(messages):
        # max_tokens=2000,
        try:
            completion = await openai.ChatCompletion.acreate(
                model=Naldia.api_model_memorize,
                messages=messages,
                temperature=1,
                top_p=0.95,
                frequency_penalty=0.0,
                presence_penalty=0.0,
            )
            response = completion.choices[0].message.content
            
        except openai.OpenAIError as e:
            response = f"Une erreur OpenAI s'est produite : {e}"
            print(response)
        return response

    def Embedding(text, model="text-embedding-ada-002"):
        text = text.replace("\n", " ")
        if Naldia.async_mode:
            return Naldia.OpenAI_embedding_async(text,model=model)
        else:
            return Naldia.OpenAI_embedding(text,model=model)
        
    async def OpenAI_embedding_async(text,model="text-embedding-ada-002"):
        return await openai.Embedding.create(input=[text], model=model)["data"][0][
            "embedding"
        ]

    def OpenAI_embedding(text,model="text-embedding-ada-002"):
        return openai.Embedding.create(input=[text], model=model)["data"][0][
            "embedding"
        ]

    #OPENAI API CALLS  
    def curl(url,headers):
        # Effectuer la requête GET
        response = requests.get(url, headers=headers)
        # Vérifier le code de réponse HTTP
        if response.status_code == 200:
            # La requête a réussi, vous pouvez traiter la réponse ici
            return response.json()
        else:
            # La requête a échoué, affichez le code de statut et le contenu de la réponse
            return (f"Code de statut HTTP : {response.status_code} \nContenu de la réponse : {response.text} ")

    def OpenAI_getOrganizationInfo():
        # URL de l'API OpenAI
        url = f'https://api.openai.com/v1/organizations/{openai.organization}'
        # En-tête d'autorisation
        headers = {
            'Authorization': f'Bearer {openai.api_key}'
        }
        print(Naldia.curl(url=url,headers=headers))

    def OpenAI_getUserInfo():
        url = f"https://api.openai.com/v1/usage?date=2023-07-10&user_public_id=$user"
        headers = {
            'Authorization': f'Bearer {openai.api_key}',
            "OpenAI-Organization": openai.organization
        }
        print(Naldia.curl(url=url,headers=headers))

        

class bcolors:
    HEADER = "\033[95m"
    OKBLUE = "\033[94m"
    OKCYAN = "\033[96m"
    OKGREEN = "\033[92m"
    WARNING = "\033[93m"
    FAIL = "\033[91m"
    ENDC = "\033[0m"
    BOLD = "\033[1m"
    UNDERLINE = "\033[4m"


def printc(str, color=bcolors.WARNING):
    print(f"{color}{str}{bcolors.ENDC}")

def inputc(str, color=bcolors.OKBLUE):
    return input(f"{color}{str}{bcolors.ENDC}")

def txt2sha(text):
    return hashlib.sha256(text.strip().encode("UTF-8")).hexdigest()





class Memorizer:
    def __init__(self, path_to_doc=None):
        self.text_path = "" # chemin vers le fichier original à mémoriser
        self.memory_path = "" # chemin vers le fichier de mémorisation (.json) 
        if path_to_doc is not None:
            self.setfile(path_to_doc)
            
    
    # charge le fichier d'origine et en déduit l'adresse de la mémoire'
    def setfile(self, path_to_doc=None):
        if path_to_doc is None:
            self.exist_info = False # pas d'info => on ne peut rien faire 
            return False
        if os.path.isfile(path_to_doc):
            # Fichier existe
            self.text_path = path_to_doc
            self.memory_path = os.path.join(os.path.dirname(self.text_path), os.path.basename(self.text_path)+".json")
            self.__build__()
        else:
            self.exist_info = False # pas d'info => on ne peut rien faire 
            return False

    def get_filename(self):
        return os.path.basename(self.text_path)
    # charge la mémoire directement sans passer par le fichier d'origine
    def load(self,memory_path):
        self.memory_path = memory_path
        self.text_path = memory_path.rstrip('.json')
        self.__build__()


    def __build__(self):
        if self.text_path != None and  not os.path.exists(self.text_path):# le document original n'existe plus 
            self.title = os.path.basename(self.text_path)
            self.info = self.get_info() # on charge la mémoire du document
            self.text = self.concatenate(self.info,"text") # on lit le texte de la mémoire du document
            self.exist_info = True
            self.calc_stats()
            return True

        if os.path.exists(self.memory_path):# la mémoire du fichier existe
            self.title = os.path.basename(self.text_path)
            self.info = self.get_info() # on charge la mémoire du document
            text_doc = self.get_text() # on lit le texte du document 
            text_info = self.concatenate(self.info,"text") # on lit le texte de la mémoire du document
            text_info = text_info.strip().replace("\n", "")
            text_doc = text_doc.strip().replace("\n", "")
            if txt2sha(text_doc) == txt2sha(text_info): # on teste si le texte est identique entre le fichier et sa mémoire
                # tout va bien, on peut commencer à chatter 
                print(f"Detected cached memories for {self.title}")
                self.exist_info = True
                self.text = text_info
                self.calc_stats()
            else:
                print(f"Detected cached memories apply on different text for {self.title}, should be Memorized again ")
                # il faut re-générer la mémoire du fichier
                self.exist_info = False
                self.text = text_info # le texte doit être celui du document, pas celui de la mémoire existante 
                # self.memorize() # on mémorise le document
                # self.info = self.get_info() # on recharge la mémoire d'aprés le nouveau texte du document 
                self.calc_stats()

        else:  
            
            if os.path.exists(self.text_path):# le textes du fichier existe
                self.exist_info = False
                self.title = os.path.basename(self.text_path)
                self.text = self.get_text().strip().replace("\n", "")
                self.calc_stats()
            else:
                self.exist_info = False
                self.title = "unknown.unk"
                self.text = "" # texte du documnet 
                self.info = None # objet json de l'objet "Memorizer"
                print(f"No cached memories for {self.title}")
                self.calc_stats()
        
        return True


    # récupération du texte du fichier, document ou url
    def get_text(self):
        self.title = os.path.basename(self.text_path)
        self.url = self.text_path
        self.suffix = os.path.splitext(self.text_path)[-1]
        if validators.url(self.url):
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36",
            }
            response = requests.get(self.url, headers=headers)
            if response.status_code == 200:
                soup = BeautifulSoup(response.content, "html.parser")
                text = soup.get_text()
            else:
                raise ValueError(f"Invalid URL! Status code {response.status_code}.")
        elif self.suffix == ".pdf":
            full_text = ""
            num_pages = 0
            with fitz.open(self.text_path) as doc:
                for page in doc:
                    num_pages += 1
                    text = page.get_text()
                    full_text += text + "\n"
            text = f"Le document fait {num_pages} page.\n" + full_text
            # utiliser une librairie "native"
            
        elif ".doc" in self.suffix:
            doc = docx.Document(self.text_path)
            fullText = []
            for para in doc.paragraphs:
                fullText.append(para.text)
            text = "\n".join(fullText)
        elif self.suffix == ".txt":
            with open(self.text_path, "r", encoding="utf8") as f:
                lines = f.readlines()
            text = "\n".join(lines)
        elif ".ppt" in self.suffix:
            text = ""
            presentation = Presentation(self.text_path)

            # Initialiser une chaîne de caractères pour stocker le texte extrait
            text = ""
            # Parcourir les slides et extraire le texte
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        else:
            raise ValueError("Invalid document path!")
            
        text = " ".join(text.split())
        return text


    def calc_stats(self):
        t= TokenPrice()
        if self.exist_info:
            self.get_info()
            self.text_words = count_words(self.text)
            self.text_token = len(Naldia.tokenizer.encode(self.text))
            self.summary = self.concatenate(self.info, "summary")
            self.summary_words = count_words(self.summary)
            self.summary_token = len(Naldia.tokenizer.encode(self.summary))
            pti = t.priceIn(self.text)
            pto = t.priceOut(self.summary)
            t.set_model("ada")
            pei = t.priceIn(self.text)
            peo = t.priceIn(self.summary)
            self.price = pti + pto + pei + peo
        else:
            self.text_words = count_words(self.text)
            self.text_token = len(Naldia.tokenizer.encode(self.text))
            self.summary = ""
            self.summary_words =0
            self.summary_token = 0
            pti = t.priceIn(self.text)
            pto = t.priceOut(self.text)*0.1021654
            t.set_model("ada")
            pei = t.priceIn(self.text)
            peo = t.priceIn(self.summary)
            self.price = pti + pto + pei + peo
        
        self.statistics=f"""
{self.title.upper()}
Texte original
    nombre de mots : {self.text_words}
    nombre de token : {self.text_token}
    cout : {round(pti, 5)}
Résumé
    nombre de mots : {self.summary_words}
    nombre de token : {self.summary_token}
    cout : {round(pto, 5)}
Embeddings :
    cout pour l'original : {round(pei, 5)}
    cout pour le résumé : {round(peo, 5)}
Prix : {round(self.price, 5)}
"""
        return self



    def memorize(self, chunk_sz=700, max_memory=150):
        printc("Mémorisation du fichier "+os.path.basename(self.text_path) +"\n",bcolors.HEADER)
        
        self.info = []
        text = self.text.replace("\n", " ").split()
        # raise error if the anticipated api usage is too massive
        test_sz = len(text) / chunk_sz
        if (test_sz) >= max_memory:
            raise ValueError(f"""
{len(text)} / {chunk_sz} = {test_sz} >= {max_memory}
Processus arrêté en raison d'un coût de traitement trop important.""")
        
        for idx in tqdm(range(0, len(text), chunk_sz)):
            chunk = " ".join(text[idx : idx + chunk_sz])
            if len(Naldia.tokenizer.encode(chunk)) > chunk_sz * 3:
                print(f"Ecarté un block peu informatif. \n\n\n{chunk} \n")
                continue
            attempts = 0
            while True:
                try:
                    summary = self.get_chunk_summary(chunk)
                    embd = Naldia.Embedding(chunk)
                    summary_embd = Naldia.Embedding(summary)
                    
                    item = self.makeInfoItem(len(self.info), chunk, embd, summary, summary_embd)
                    """
                    item = {
                        "id": len(self.info),
                        "text": chunk,
                        "embd": embd,
                        "summary": summary,
                        "summary_embd": summary_embd,
                    }
                    """
                    self.info.append(item)
                    Naldia.sleep()
                    break
                except Exception as e:
                    attempts += 1
                    printc(f"\nApi call failed {attempts}, retrying ...",bcolors.WARNING)
                    if attempts >= 3:
                        raise Exception(f"{str(e)}")
                    Naldia.sleep()
        with jsonlines.open(self.memory_path, mode="w") as f:
            f.write(self.info)
            # print(f"Information enregistrée dans {self.memory_path}")
            self.exist_info = True

        self.calc_stats()
        self.save_summary()
        
    def makeInfoItem(self, id=-1,chunk="",embd=[],summary="",summary_embd=[]):
        if id ==-1:
            id = random.randint(100000, 10000000)
        item = {
            "id": id,
            "text": chunk,
            "embd": embd,
            "summary": summary,
            "summary_embd": summary_embd,
        }
        return item
    

    def get_chunk_summary(self, chunk):
        # content = "The following is a passage fragment. Please summarize what information the readers can take away from it:"
        content = "Ce qui suit est un fragment du document. Veuillez résumer les informations que les lecteurs peuvent en retirer:"
        content += "\n" + chunk
        messages = [{"role": "user", "content": content}]
        
        summary = Naldia.ask(messages)
        return summary

    # affecte l'info à l'objet mémory
    def set_info(self,info):
        self.info = info

    # retourne l'information du mémo dans un objet "info"
    def get_info(self):
        if not os.path.exists(self.memory_path):
            return self.makeInfoItem()
        with open(self.memory_path, "r", encoding="utf8") as f:
            for line in f:
                info = json.loads(line)
                self.info = info
        return info

    def save_summary(self):
        if self.exist_info:
            summary = self.concatenate(self.info, "summary")
            with open(f'{self.memory_path}.txt', 'w', encoding="utf8") as file:
                file.write(f"""
{self.statistics}
{summary}
""")
        else:
            printc("Le document doit d'abord être mémorisé", bcolors.WARNING)
    

    def get_summary(self):
        return self.concatenate(self.info, "summary")

    # concatène les attribut de l'info en mémoire
    def concatenate(self, json_data, key="summary" ):
        """item = {"id": id, "text": chunk, "embd": embd, "summary": summary, "summary_embd": summary_embd,}"""
        summaries = [item[key].strip().replace("\n", "") for item in json_data]
        
        # l'espace de concaténation est ULTRA IMPORTANT !!!!!!
        concatenated_summary = ' '.join(summaries) 
        return concatenated_summary


def count_words(text):
    word_count = len(text.split())
    return word_count


class TokenPrice():
    def __init__(self,model="gpt3"):
        self.set_model(model)

    def set_model(self,model):
        self.model = model
        if self.model == "gpt-3.5-turbo" or self.model == "gpt3" :
            self.input_cost = 0.0015 # per 1000 token
            self.ouput_cost = 0.002
            self.model = "gpt-3.5-turbo"

        elif self.model == "text-davinci-003" or self.model == "davinci" :
            self.input_cost = 0.0200 
            self.ouput_cost = 0.002
            self.model = "text-davinci-003"

        elif self.model == "text-embedding-ada-002" or self.model == "ada" :
            self.input_cost = 0.0001 
            self.ouput_cost = 0
            self.model = "text-embedding-ada-002"
        self.Tokenizer = tiktoken.encoding_for_model(self.model)
        return self

    def priceIn(self,text):
        self.nbr_token = len(self.Tokenizer.encode(text))
        if self.model == "gpt-3.5-turbo" or self.model == "gpt3" :
            self.input_price =  self.nbr_token * self.input_cost / 1000
        elif self.model == "text-embedding-ada-002" or self.model == "ada" :
            self.input_price = len(self.Tokenizer.encode(text)) * self.input_cost / 1000
        return self.input_price
    
    def priceOut(self,text):
        self.nbr_token = len(self.Tokenizer.encode(text))
        if self.model == "gpt-3.5-turbo" or self.model == "gpt3" :
            self.ouput_price =  self.nbr_token * self.ouput_cost / 1000
        elif self.model == "text-embedding-ada-002" or self.model == "ada" :
            self.ouput_price = len(self.Tokenizer.encode(text)) * self.ouput_cost / 1000
        return self.ouput_price
    


class Chatter:
    can_chat = False
    conversation=[]
    def __init__(self, *memos: Memorizer):
        if memos:
            if len(memos)==1:
                printc("Un seul mémo chargé dans la discussion",bcolors.HEADER)
                self.load(memos[0])
            else:
                printc("plusieurs mémos chargés dans la discussion",bcolors.HEADER)
                self.memo = Memorizer() 
                combined_info = []
                for memo in memos:
                    if memo.info == None :
                        printc(f"Le mémo {memo.title} n'a pas été Mémorisé, il est écarté de la discution",bcolors.WARNING)
                    else :
                        info = memo.get_info()
                        combined_info.extend(info)
                if len(combined_info) <1:
                    printc(f"Aucun mémo n'est mémorisé",bcolors.WARNING)
                else:
                    self.memo.set_info(combined_info)
                    self.info = combined_info
                    self.can_chat = True
        # Traitement lorsque qu'aucun mémo n'est fourni
        else:
            printc("Aucun mémo n'a été fourni",bcolors.WARNING)

    # permet de charger un mémo aprés avoir initialisé une instance de l'objet Chatter
    def load(self, memo: Memorizer):
        if memo.info == None :
            print(f"Le mémo {memo.title} n'a pas été Mémorisé ... chat impossible")
            self.can_chat = False
        else:
            self.memo = memo
            self.info = self.memo.get_info()
            self.can_chat = True

    def ask(self, q):
        if not self.can_chat:
            return "Chat impossible ..."
        attempts = 0
        if len(Naldia.tokenizer.encode(q)) > 200:
            return "la question est trop longue"
        while True:
            try:
                self.conversation.append({"role":"user","content":q})
                response = self.answer(q)
                self.conversation.append({"role":"system","content":response})
                self.output(response)
                # Naldia.sleep()
                break
            except Exception as e:
                attempts += 1
                if attempts >= 3:
                    raise Exception(f"{str(e)}")
                # Naldia.sleep()
        return response

    def answer(self, q):
        if Naldia.offline:
            answer = "Naldia est hors ligne, réponse standard."
        else:
            # Naldia.api_model_current = Naldia.api_model_memorize
            q_embd = Naldia.OpenAI_embedding(q, model="text-embedding-ada-002")
            retrieved_indices = self.retrieve(q_embd)
            answer = self.generate_answer(q, retrieved_indices)
        return answer

    # return the indices of top three related texts
    def retrieve(self, q_embd):
        text_embds = []
        summary_embds = []
        for item in self.info:
            text_embds.append(item["embd"])
            summary_embds.append(item["summary_embd"])
        # compute the cos sim between info_embds and q_embd
        text_cos_sims = np.dot(text_embds, q_embd) / (
            norm(text_embds, axis=1) * norm(q_embd)
        )
        summary_cos_sims = np.dot(summary_embds, q_embd) / (
            norm(summary_embds, axis=1) * norm(q_embd)
        )
        cos_sims = text_cos_sims + summary_cos_sims
        top_args = np.argsort(cos_sims).tolist()
        top_args.reverse()
        indices = top_args[0:3]
        return indices

    def generate_answer(self, q, retrieved_indices):
        while True:
            sorted_indices = sorted(retrieved_indices)
            retrieved_text = [self.info[idx]["text"] for idx in sorted_indices]
            content = self.get_qa_content(q, retrieved_text)
            if len(Naldia.tokenizer.encode(content)) > 3800:
                retrieved_indices = retrieved_indices[:-1]
                print("Reflexion en cours...")
                if not retrieved_indices:
                    raise ValueError("Failed to respond.")
            else:
                break
        
        messages = [{"role": "user", "content": content}]
        answer = Naldia.ask(messages)
        return answer

    def get_qa_content1(self, q, retrieved_text):
        # content = "After reading some relevant passage fragments from the same document, please respond to the following query. Note that there may be typographical errors in the passages due to the text being fetched from a PDF file or web page."
        # content += "\nAvoid explicitly using terms such as 'passage 1, 2 or 3' in your answer as the questioner may not know how the fragments are retrieved. Please use the same language as in the query to respond."
        content = """Après avoir lu quelques passages pertinents d'un même document, veuillez répondre à la Question suivante.
        Notez qu'il peut y avoir des erreurs typographiques dans les passages.
        """
        content += "\nQuestion: " + q
        for i in range(len(retrieved_text)):
            content += "\npassage " + str(i + 1) + ": " + retrieved_text[i]
        content += """
        Évitez d'utiliser explicitement des termes tels que "passage 1, 2 ou 3"
        """
        return content
    
    def get_qa_content(self, q, retrieved_text):

        info = self.RelevantInfo(retrieved_text)
        content = f"""
Aprés avoir lu quelques passages pertinents d'un même document, répondre à la question :
"{q}"
Utilisez vos connaissances si certaines informations ne sont pas explicites dans le texte pour enrichir votre réponse.
Répondre dans la même langue que la question.
Mettre en forme votre réponse au format Markdown.

Texte : {info}
        """
        # print (content)
        return content

    def RelevantInfo(self, retrieved_text):
        text =""
        for i in range(len(retrieved_text)):
            text += "\npassage " + str(i + 1) + ": " + retrieved_text[i]
        return text
    
    def save(self):
        with open("conversation.txt", "a+") as fichier:
            for element in self.conversation:
                fichier.write(element + "\n")




    def chat(self,q=""):
        # initialyse chatted memo and extract its informations
        if not self.can_chat:
            return "Chat impossible ..."
        chatting = True
        while chatting:
            q = self.get_question()
            if len(Naldia.tokenizer.encode(q)) > 200:
                raise ValueError("Input query is too long!")
            elif q == "":
                chatting = True
            elif q.startswith("#"):
                self.special_prompt(q[1:])
            elif q == "#quit" or q == "#exit" :
                chatting = False
            else:
                attempts = 0
                while True:
                    try:
                        self.conversation.append({"role":"user","content":q})
                        response = self.answer(q)
                        self.conversation.append({"role":"system","content":response})
                        self.output(response)
                        # Naldia.sleep()
                        break
                    except Exception as e:
                        attempts += 1
                        if attempts >= 3:
                            raise Exception(f"{str(e)}")
                        # Naldia.sleep()

    # entrée utilisateur pour poser uen question
    def get_question(self):
        q = input(f"Entrer votre question ou #help (offline : {Naldia.offline}): ")
        return q
    def output(self,response):
        if Naldia.api_model_current == Naldia.api_model_ask:
            color = bcolors.OKBLUE
        else:
            color = bcolors.OKGREEN
        print("",f"{color}{response}{bcolors.ENDC}","")
    
    def special_prompt(self,q):
        if q=="help":
            printc("""
    HELP : using spécial #keywords
    #quit or #exit : exit conversation 
    #save : save the current conversation 
    #info : displays Naldia settings
    #online : set the chat online
    #offline : set the chat offline (openAI API will not be contacted, for debug and test purposes)
    #api_model_ask : set Naldia model to "davinci"
    #api_model_memorize : set Naldia model to "gpt3"
    #api_model_current : displays the current model

    #get_organization : retrieve organization information

    """,bcolors.OKBLUE)

        if q=="info":
            printc(f"""
Naldia.offline : {Naldia.offline}
Naldia.api_call_wait : {Naldia.api_call_wait}
Naldia.api_model_current : {Naldia.api_model_current}
Chatter.can_chat : {self.can_chat}

    """)
            print(Naldia.offline)

        if q=="online":
            Naldia.offline=False
            print(Naldia.offline)

        if q=="offline":
            Naldia.offline=True
            print(Naldia.offline)

        if q=="api_model_memorize":
            Naldia.set_model(Naldia.api_model_memorize)
            print(Naldia.api_model_current)

        if q=="api_model_ask":
            Naldia.set_model(Naldia.api_model_ask)
            print(Naldia.api_model_current)

        if q=="api_model_current":
            print(Naldia.api_model_current)

        if q=="save":
            self.save()
        
        if q=="get_organization":
            print(Naldia.OpenAI_getOrganizationInfo())